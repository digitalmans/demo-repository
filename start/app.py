#!/usr/bin/env python
# -*- coding:utf-8 -*-

"""
用户认证系统 - Flask后端
提供登录和注册功能（使用MySQL数据库）
"""

import os
import sys
import json

# 在导入任何模块之前设置 TensorFlow 环境变量，避免导入时的警告和错误
os.environ["TF_CPP_MIN_LOG_LEVEL"] = "3"  # 屏蔽 TensorFlow 的 INFO、WARNING、ERROR 日志
os.environ["TF_ENABLE_ONEDNN_OPTS"] = "0"  # 禁用 oneDNN 优化，避免警告
from flask import Flask, render_template, request, jsonify, session, redirect, url_for, send_file
from functools import wraps
from werkzeug.utils import secure_filename
from config import SECRET_KEY
from database import (
    init_database, register_user, verify_user, get_user_info,
    get_all_users, update_user_role, delete_user,

    save_qa_robot_history, get_qa_robot_history, get_qa_robot_history_count,
    delete_qa_robot_history,
    get_tts_audio_cache, save_tts_audio_cache,
    save_qa_discussion, get_qa_discussions, get_qa_discussion_count,
    delete_qa_discussion,
    save_qa_comment, get_qa_comments, delete_qa_comment,
    # 增强功能
    get_qa_discussions_extended, get_qa_discussion_count_extended,
    get_qa_discussion_detail, vote_qa_target, get_user_votes,
    get_qa_comments_extended, toggle_favorite, is_favorited,
    get_user_favorited_ids, get_user_profile,
    get_user_discussions, get_user_comments, get_user_favorites,
    update_user_profile
)

# 添加 voice_assistant 目录到路径，以便导入 asr 和 tts 模块
voice_assistant_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'voice_assistant')
if voice_assistant_path not in sys.path:
    sys.path.insert(0, voice_assistant_path)

# 移除了电影问答目录到路径


# 添加 ask_answer_robot 目录到路径，以便导入问答机器人服务
qa_robot_path = os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 
                              'ask_answer_robot')
if qa_robot_path not in sys.path:
    sys.path.insert(0, qa_robot_path)

# 添加 voice_assistant 目录到路径，以便导入翻译模块
if voice_assistant_path not in sys.path:
    sys.path.insert(0, voice_assistant_path)

try:
    from asr import asr, DemoError as ASRError
    from tts import tts, DemoError as TTSError
    VOICE_ASSISTANT_AVAILABLE = True
except ImportError as e:
    print(f"警告: 无法导入语音助手模块: {e}")
    VOICE_ASSISTANT_AVAILABLE = False

# 导入翻译模块
try:
    from translator import translate, translate_to_chinese, translate_to_english, translate_to_japanese, TranslationError
    TRANSLATOR_AVAILABLE = True
except ImportError as e:
    print(f"警告: 无法导入翻译模块: {e}")
    TRANSLATOR_AVAILABLE = False

app = Flask(__name__)
app.secret_key = SECRET_KEY

@app.context_processor
def inject_view_variables():
    try:
        deepseek_configured = DEEPSEEK_AVAILABLE and deepseek_service and deepseek_service.is_configured()
    except Exception:
        deepseek_configured = False
    return {
        'deepseek_presets': DEEPSEEK_PRESETS if 'DEEPSEEK_PRESETS' in globals() else {},
        'default_deepseek_preset': DEFAULT_DEEPSEEK_PRESET if 'DEFAULT_DEEPSEEK_PRESET' in globals() else '',
        'deepseek_configured': deepseek_configured,
        'qa_robot_available': QA_ROBOT_AVAILABLE if 'QA_ROBOT_AVAILABLE' in globals() else False
    }


# 导入问答机器人服务
QA_ROBOT_AVAILABLE = False
qa_robot_service = None
try:
    from qa_service import QAService
    QA_ROBOT_AVAILABLE = True
    print("✓ 问答机器人模块导入成功")
except ImportError as e:
    print(f"✗ 警告: 无法导入问答机器人模块: {e}")
    print(f"  提示: 请确保已安装 py2neo 和 jieba 等依赖")
    QA_ROBOT_AVAILABLE = False

# 初始化问答机器人服务
if QA_ROBOT_AVAILABLE:
    try:
        print("正在初始化问答机器人服务...")
        # 默认不使用BERT，限制初始加载数量，避免启动时等待过久
        qa_robot_service = QAService(use_bert=False, max_initial_load=20000)
        print("✓ 问答机器人服务初始化成功")
        print("  提示: 当前使用text_similarity方法（BM25+Jaccard+编辑距离）")
        print("  提示: 初始仅加载20000条数据，如需更多请修改 max_initial_load 参数")
        print("  如需使用BERT，请修改代码设置 use_bert=True 并启动BERT服务")
    except Exception as e:
        print(f"✗ 警告: 问答机器人服务初始化失败: {e}")
        import traceback
        traceback.print_exc()
        print("\n提示: 请确保:")
        print("  1. 已安装 py2neo: pip install py2neo")
        print("  2. 已安装 jieba: pip install jieba")
        print("  3. Neo4j 服务已启动并可以连接")
        print("  4. 已导入语料库数据到Neo4j")
        QA_ROBOT_AVAILABLE = False
        qa_robot_service = None

# 导入DeepSeek服务
DEEPSEEK_AVAILABLE = False
deepseek_service = None
DEEPSEEK_PRESETS = {}
try:
    from deepseek_service import (
        DeepSeekService,
        DeepSeekAPIError,
        DEEPSEEK_PRESETS,
        DEFAULT_PRESET as DEFAULT_DEEPSEEK_PRESET,
    )
    deepseek_service = DeepSeekService()
    DEEPSEEK_AVAILABLE = True
    if deepseek_service.is_configured():
        print("✓ DeepSeek API 已配置")
    else:
        print("提示: 未配置 DEEPSEEK_API_KEY，DeepSeek回答将回退到本地知识库")
except ImportError as e:
    print(f"警告: 无法导入DeepSeek服务模块: {e}")
    DeepSeekAPIError = Exception
    DEFAULT_DEEPSEEK_PRESET = "deepseek_qa"

# 语音助手配置
if VOICE_ASSISTANT_AVAILABLE:
    app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024  # 限制上传文件大小为16MB
    app.config['UPLOAD_FOLDER'] = os.path.join(voice_assistant_path, 'uploads')
    app.config['OUTPUT_FOLDER'] = os.path.join(voice_assistant_path, 'outputs')
    
    # 确保上传和输出目录存在
    os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
    os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)
    
    # 允许的音频文件扩展名
    ALLOWED_EXTENSIONS = {'pcm', 'wav', 'amr', 'mp3', 'm4a'}
    
    def allowed_file(filename):
        """检查文件扩展名是否允许"""
        return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS


def login_required(f):
    """登录验证装饰器"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'username' not in session:
            return redirect(url_for('login'))
        return f(*args, **kwargs)
    return decorated_function


def admin_required(f):
    """管理员权限验证装饰器"""
    @wraps(f)
    def decorated_function(*args, **kwargs):
        if 'username' not in session:
            return redirect(url_for('login'))
        if session.get('role') != 'admin':
            return jsonify({'success': False, 'error': '需要管理员权限'}), 403
        return f(*args, **kwargs)
    return decorated_function


@app.route('/')
def index():
    """首页，重定向到登录页"""
    if 'username' in session:
        return redirect(url_for('qa_robot'))
    return redirect(url_for('login'))


@app.route('/login', methods=['GET', 'POST'])
def login():
    """登录页面"""
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not password:
            return jsonify({'success': False, 'error': '用户名和密码不能为空'}), 400
        
        # 验证用户
        success, message, user_data = verify_user(username, password)
        
        if not success:
            return jsonify({'success': False, 'error': message}), 401
        
        # 登录成功，设置session
        session['username'] = username
        session['user_id'] = user_data['id']
        session['role'] = user_data.get('role', 'user')
        
        # 根据角色跳转
        if user_data.get('role') == 'admin':
            return jsonify({
                'success': True, 
                'message': '管理员登录成功', 
                'user': user_data,
                'redirect': '/admin'
            })
        else:
            return jsonify({
                'success': True, 
                'message': '登录成功', 
                'user': user_data,
                'redirect': '/voice_assistant'
            })
    
    # GET请求，返回登录页面
    if 'username' in session:
        return redirect(url_for('qa_robot'))
    return render_template('login.html')


@app.route('/register', methods=['GET', 'POST'])
def register():
    """注册页面"""
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        confirm_password = data.get('confirm_password', '').strip()
        email = data.get('email', '').strip() or None
        
        # 验证输入
        if not username or not password:
            return jsonify({'success': False, 'error': '用户名和密码不能为空'}), 400
        
        if len(username) < 3:
            return jsonify({'success': False, 'error': '用户名至少需要3个字符'}), 400
        
        if len(password) < 6:
            return jsonify({'success': False, 'error': '密码至少需要6个字符'}), 400
        
        if password != confirm_password:
            return jsonify({'success': False, 'error': '两次输入的密码不一致'}), 400
        
        # 注册新用户
        success, message = register_user(username, password, email)
        
        if not success:
            status_code = 409 if '已存在' in message else 400
            return jsonify({'success': False, 'error': message}), status_code
        
        return jsonify({'success': True, 'message': '注册成功，请登录'})
    
    # GET请求，返回注册页面
    return render_template('register.html')


@app.route('/logout')
def logout():
    """登出"""
    session.pop('username', None)
    session.pop('user_id', None)
    session.pop('role', None)
    return redirect(url_for('login'))


@app.route('/voice_assistant')
@login_required
def voice_assistant():
    """语音助手页面（重定向到问答机器人）"""
    # 默认跳转到问答机器人
    return redirect(url_for('qa_robot'))


@app.route('/qa_robot')
@login_required
def qa_robot():
    """问答机器人页面（需要登录）"""
    deepseek_configured = DEEPSEEK_AVAILABLE and deepseek_service and deepseek_service.is_configured()
    if not QA_ROBOT_AVAILABLE and not DEEPSEEK_AVAILABLE:
        return render_template('voice_assistant.html', 
                             username=session.get('username'),
                             current_view='qa_robot',
                             error='问答机器人模块不可用，且DeepSeek服务模块不可用')
    return render_template('voice_assistant.html', 
                         username=session.get('username'),
                         current_view='qa_robot',
                         deepseek_presets=DEEPSEEK_PRESETS,
                         default_deepseek_preset=DEFAULT_DEEPSEEK_PRESET,
                         deepseek_configured=deepseek_configured,
                         qa_robot_available=QA_ROBOT_AVAILABLE)


@app.route('/mood_pixel')
@login_required
def mood_pixel():
    """心情拼豆页面（需要登录）"""
    return render_template('voice_assistant.html',
                           username=session.get('username'),
                           current_view='mood_pixel')


@app.route('/api/mood_pixel/text_to_image', methods=['POST'])
@login_required
def mood_pixel_text_to_image():
    """使用 Moark FLUX-1-schnell 进行文生图"""
    data = request.get_json() or {}
    prompt = data.get('prompt', '').strip()
    size = data.get('size', '1024x1024')
    if not prompt:
        return jsonify({'error': 'Prompt is required'}), 400

    try:
        from openai import OpenAI
        import base64
        import requests as req

        client = OpenAI(
            base_url="https://api.moark.com/v1",
            api_key="QQFL0VLH1MMPVEOASHZAOTMJOCTXC2XHD4MWBO1Q",
        )

        response = client.images.generate(
            prompt=prompt,
            model="flux-1-schnell",
            size=size,
            extra_body={
                "guidance_scale": 7.5,
                "seed": 42,
                "lora_weights": [],
                "lora_scale": 0,
                "num_images_per_prompt": 1,
                "width": 0,
                "height": 0,
                "negative_prompt": "blurry, low quality, distorted",
            },
            response_format="b64_json",
        )

        for i, image_data in enumerate(response.data):
            if getattr(image_data, 'b64_json', None):
                return jsonify({
                    'success': True,
                    'image_base64': 'data:image/jpeg;base64,' + image_data.b64_json
                })
            elif getattr(image_data, 'url', None):
                return jsonify({
                    'success': True,
                    'image_url': image_data.url
                })
        return jsonify({'error': 'No image data returned'}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500


@app.route('/api/mood_pixel/image_to_image', methods=['POST'])
@login_required
def mood_pixel_image_to_image():
    """使用 Moark FLUX.1-Kontext-dev 进行图生图"""
    if 'image' not in request.files:
        return jsonify({'error': 'Image file is required'}), 400

    file = request.files['image']
    prompt = request.form.get('prompt', 'Convert to colorful clean perler bead pixel art pattern')

    try:
        import requests as req
        from requests_toolbelt import MultipartEncoder
        import os
        import tempfile

        suffix = os.path.splitext(file.filename or '.jpg')[1] or '.jpg'
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name

        API_URL = "https://api.moark.com/v1/images/edits"
        API_TOKEN = "QQFL0VLH1MMPVEOASHZAOTMJOCTXC2XHD4MWBO1Q"
        headers = {
            "Authorization": f"Bearer {API_TOKEN}"
        }

        fields = [
            ("prompt", prompt),
            ("model", "FLUX.1-Kontext-dev"),
            ("size", "1024x1024"),
            ("steps", "20"),
            ("guidance_scale", "2.5"),
            ("seed", "42"),
            ("return_image_quality", "80"),
            ("return_image_format", "b64_json"),
            ("lora_scale", "0"),
            ("width", "0"),
            ("height", "0"),
        ]

        mime_type = file.mimetype or "image/jpeg"
        with open(tmp_path, "rb") as f:
            fields.append(("image", (os.path.basename(tmp_path), f.read(), mime_type)))

        encoder = MultipartEncoder(fields)
        headers["Content-Type"] = encoder.content_type

        response = req.post(API_URL, headers=headers, data=encoder, timeout=60)
        try:
            os.unlink(tmp_path)
        except Exception:
            pass

        result = response.json()
        if "data" in result and len(result["data"]) > 0:
            img_data = result["data"][0]
            if "b64_json" in img_data:
                return jsonify({
                    'success': True,
                    'image_base64': 'data:image/jpeg;base64,' + img_data["b64_json"]
                })
            elif "url" in img_data:
                return jsonify({
                    'success': True,
                    'image_url': img_data["url"]
                })
        return jsonify({'error': 'Generation failed or returned unexpected format', 'details': str(result)}), 500
    except Exception as e:
        return jsonify({'error': str(e)}), 500



@app.route('/qa_discussion')
@login_required
def qa_discussion():
    """QA讨论区页面（需要登录）"""
    return render_template('voice_assistant.html',
                         username=session.get('username'),
                         user_id=session.get('user_id'),
                         current_view='qa_discussion')


@app.route('/digital_human')
@login_required
def digital_human():
    """3D 数字人页面（直达 3D 数字人系统）"""
    return redirect("http://127.0.0.1:5173")



@app.route('/feature_experience')
@login_required
def feature_experience():
    """功能体验页面（需要登录）"""
    return render_template('voice_assistant.html', 
                         username=session.get('username'),
                         current_view='feature_experience')


@app.route('/settings')
@login_required
def settings():
    """系统设置页面（需要登录）"""
    return render_template('voice_assistant.html', 
                         username=session.get('username'),
                         current_view='settings')


# ==================== SPA 局部视图 API ====================
# 这些路由为前端 SPA 导航提供 HTML 片段，无需整页刷新

@app.route('/api/partial/qa_robot')
@login_required
def partial_qa_robot():
    """返回问答机器人 HTML 片段"""
    deepseek_configured = DEEPSEEK_AVAILABLE and deepseek_service and deepseek_service.is_configured()
    return render_template('qa_robot_content.html',
                           username=session.get('username'),
                           current_view='qa_robot',
                           deepseek_presets=DEEPSEEK_PRESETS,
                           default_deepseek_preset=DEFAULT_DEEPSEEK_PRESET,
                           deepseek_configured=deepseek_configured,
                           qa_robot_available=QA_ROBOT_AVAILABLE)

@app.route('/api/partial/mood_pixel')
@login_required
def partial_mood_pixel():
    """返回心情拼豆 HTML 片段"""
    return render_template('mood_pixel_content.html',
                           username=session.get('username'),
                           current_view='mood_pixel')


@app.route('/api/partial/qa_discussion')
@login_required
def partial_qa_discussion():
    """返回讨论区 HTML 片段"""
    return render_template('qa_discussion_content.html',
                           username=session.get('username'),
                           user_id=session.get('user_id'),
                           current_view='qa_discussion')

@app.route('/api/partial/feature_experience')
@login_required
def partial_feature_experience():
    """返回功能体验 HTML 片段"""
    return render_template('feature_experience_content.html',
                           username=session.get('username'),
                           current_view='feature_experience')

@app.route('/api/partial/settings')
@login_required
def partial_settings():
    """返回系统设置 HTML 片段"""
    return render_template('settings_content.html',
                           username=session.get('username'),
                           current_view='settings')

@app.route('/api/partial/user_profile')
@login_required
def partial_user_profile():
    """返回个人主页 HTML 片段"""
    username = request.args.get('username', session.get('username'))
    return render_template('user_profile.html',
                           username=session.get('username'),
                           profile_username=username,
                           current_view='user_profile')


# 快速问答对映射（直接返回，不进行检索）
QUICK_QA_MAP = {
    "太阳系中体积最大的行星是哪一颗": "木星",
    "太阳系中体积最大的行星是哪一颗？": "木星",
    "太阳系体积最大的行星": "木星",
    "体积最大的行星": "木星",
    "人体最大的器官是什么": "皮肤",
    "人体最大的器官是什么？": "皮肤",
    "人体最大器官": "皮肤",
    "最大的器官": "皮肤",
    "澳大利亚的首都是悉尼吗": "不是，是堪培拉",
    "澳大利亚的首都是悉尼吗？": "不是，是堪培拉",
    "澳大利亚首都是悉尼吗": "不是，是堪培拉",
    "澳洲首都是悉尼吗": "不是，是堪培拉",
    "蒙娜丽莎的作者是谁": "达·芬奇",
    "蒙娜丽莎的作者是谁？": "达·芬奇",
    "蒙娜丽莎作者": "达·芬奇",
    "水的化学式是什么": "H2O",
    "水的化学式是什么？": "H2O",
    "水的化学式": "H2O",
    "企鹅主要生活在北极还是南极": "南极",
    "企鹅主要生活在北极还是南极？": "企鹅主要生活在南极",
    "企鹅生活在哪": "南极",
    "企鹅在北极还是南极": "南极",
    "光在真空中的传播速度大约是多少": "约 30 万公里/秒",
    "光在真空中的传播速度大约是多少？": "约 30 万公里/秒",
    "光速是多少": "约 30 万公里/秒",
    "光的速度": "约 30 万公里/秒",
    "著名的相对论是谁提出的": "爱因斯坦",
    "著名的相对论是谁提出的？": "爱因斯坦",
    "相对论是谁提出的": "爱因斯坦",
    "相对论的提出者": "爱因斯坦",
    "一场标准的足球比赛中，每队有几名球员上场": "11名",
    "一场标准的足球比赛中，每队有几名球员上场？": "11名",
    "足球比赛每队几人": "11名",
    "足球每队几个人": "11名",
    "世界上最高的山峰是哪座": "珠穆朗玛峰",
    "世界上最高的山峰是哪座？": "珠穆朗玛峰",
    "世界最高峰": "珠穆朗玛峰",
    "最高的山峰": "珠穆朗玛峰"
}

def match_quick_qa(question):
    """
    匹配快速问答对（智能匹配，支持多种问法）
    :param question: 用户问题
    :return: 如果匹配，返回答案；否则返回None
    """
    import re
    
    # 去除标点符号和空格，进行标准化匹配
    normalized_question = re.sub(r'[，。！？、；：\s]', '', question.strip())
    
    # 1. 直接匹配
    if question in QUICK_QA_MAP:
        return QUICK_QA_MAP[question]
    
    # 2. 标准化匹配（去除标点后匹配）
    for key, value in QUICK_QA_MAP.items():
        normalized_key = re.sub(r'[，。！？、；：\s]', '', key.strip())
        if normalized_question == normalized_key:
            return value
    
    # 3. 智能关键词匹配（处理各种问法变体）
    question_lower = question.lower()
    
    # 太阳系体积最大的行星
    if ("太阳系" in question or "太阳" in question) and ("体积最大" in question or "最大" in question) and ("行星" in question or "星球" in question):
        return "木星"
    
    # 人体最大器官
    if ("人体" in question or "人" in question) and "最大" in question and ("器官" in question or "部位" in question):
        return "皮肤"
    
    # 澳大利亚首都
    if ("澳大利亚" in question or "澳洲" in question) and "首都" in question:
        if "悉尼" in question:
            return "不是，是堪培拉"
        return "堪培拉"
    
    # 蒙娜丽莎作者
    if "蒙娜丽莎" in question and ("作者" in question or "谁" in question or "画" in question):
        return "达·芬奇"
    
    # 水的化学式
    if "水" in question and ("化学式" in question or "化学" in question or "分子式" in question):
        return "H2O"
    
    # 企鹅生活在哪里
    if "企鹅" in question:
        if "北极" in question:
            return "企鹅主要生活在南极，不是北极"
        if "南极" in question:
            return "南极"
        if "哪里" in question or "哪" in question:
            return "南极"
    
    # 光速
    if "光" in question and ("速度" in question or "传播" in question or "光速" in question):
        if "真空" in question or "真空中" in question:
            return "约 30 万公里/秒"
        return "约 30 万公里/秒"
    
    # 相对论
    if "相对论" in question and ("提出" in question or "谁" in question or "发明" in question):
        return "爱因斯坦"
    
    # 足球比赛人数
    if "足球" in question and ("球员" in question or "队员" in question or "人数" in question or "几个人" in question or "多少" in question):
        return "11名"
    
    # 世界最高峰
    if ("最高" in question or "最高峰" in question) and ("山" in question or "峰" in question):
        if "世界" in question or "全球" in question:
            return "珠穆朗玛峰"
        return "珠穆朗玛峰"
    
    return None


@app.route('/api/qa_robot/ask', methods=['POST'])
@login_required
def api_qa_robot_ask():
    """问答机器人API"""
    try:
        data = request.get_json() or {}
        question = str(data.get('question') or '').strip()
        preset = str(data.get('preset') or 'Qwen3-32B').strip() or 'Qwen3-32B'
        deepseek_api_key = str(data.get('deepseek_api_key') or '').strip()
        use_local_only = preset == 'local'
        force_deepseek = (data.get('force_deepseek') == True) or (preset == 'deepseek_only')
        
        if not question:
            return jsonify({'success': False, 'error': '问题不能为空'}), 400

        request_deepseek_service = deepseek_service
        if deepseek_api_key and DEEPSEEK_AVAILABLE:
            request_deepseek_service = DeepSeekService(api_key=deepseek_api_key)
        elif DEEPSEEK_AVAILABLE and request_deepseek_service and not request_deepseek_service.is_configured():
            request_deepseek_service = DeepSeekService()
            
        deepseek_configured = (
            DEEPSEEK_AVAILABLE and
            request_deepseek_service and
            request_deepseek_service.is_configured()
        )
        if not QA_ROBOT_AVAILABLE and (use_local_only or not deepseek_configured):
            return jsonify({'success': False, 'error': '问答机器人模块不可用，且当前请求无法使用DeepSeek API'}), 503
        
        use_kb = data.get('use_kb', True)
        
        # 首先检查是否是快速问答对（如果强制使用DeepSeek，则跳过快速匹配）
        quick_answer = None if force_deepseek else match_quick_qa(question)
        if quick_answer:
            # 快速问答对也保存历史记录
            similarity = 1.0  # 完全匹配，相似度为1.0
            if 'user_id' in session and 'username' in session:
                save_qa_robot_history(
                    session['user_id'],
                    session['username'],
                    question,
                    quick_answer,
                    similarity
                )
            
            return jsonify({
                'success': True,
                'question': question,
                'answer': quick_answer,
                'similarity': similarity,
                'matched_question': question,
                'source': 'quick_qa',
                'preset': 'quick_qa',
                'alternatives': []
            })
        
        # 如果不是快速问答对，并且启用了本地知识库（或者是仅本地问答模式），则做本地检索
        result = None
        if not force_deepseek and (use_local_only or use_kb) and QA_ROBOT_AVAILABLE and qa_robot_service:
            result = qa_robot_service.ask(question, top_k=3, threshold=0.1)

        answer = None
        similarity = result.get('similarity') if result else None
        source = 'local'
        notice = None
        deepseek_meta = {}

        if use_local_only:
            if not result:
                return jsonify({
                    'success': False,
                    'error': '没有找到相关答案，请切换到DeepSeek综合问答或换一种方式提问。'
                }), 404
            answer = result.get('answer', '')
        else:
            if deepseek_configured:
                try:
                    deepseek_result = request_deepseek_service.ask(question, preset_key=preset, local_result=result)
                    answer = deepseek_result['answer']
                    source = 'deepseek'
                    deepseek_meta = deepseek_result
                except DeepSeekAPIError as e:
                    if not result:
                        return jsonify({'success': False, 'error': f'DeepSeek API调用失败: {str(e)}'}), 502
                    answer = result.get('answer', '')
                    source = 'local_fallback'
                    notice = f'DeepSeek API调用失败，已回退到本地知识库: {str(e)}'
            elif result:
                answer = result.get('answer', '')
                source = 'local_fallback'
                notice = '未配置 DEEPSEEK_API_KEY，已回退到本地知识库。'
            else:
                return jsonify({
                    'success': False,
                    'error': '未配置 DEEPSEEK_API_KEY，且本地知识库没有找到相关答案。'
                }), 503

        if not answer:
            return jsonify({'success': False, 'error': '没有生成有效答案'}), 500

        # 限制答案长度，最多2000字，避免前端和语音播报过长
        MAX_LENGTH = 2000
        if answer and len(answer) > MAX_LENGTH:
            answer = answer[:MAX_LENGTH] + "\n\n(单次回答最多生成2000字)"
        
        if 'user_id' in session and 'username' in session:
            save_qa_robot_history(
                session['user_id'],
                session['username'],
                question,
                answer,
                similarity
            )

        alternatives = []
        if result:
            for alt in result.get('alternatives', []):
                if len(alt) >= 3:
                    alternatives.append({
                        'question': alt[0],
                        'answer': alt[1],
                        'similarity': round(alt[2], 3),
                        'source': alt[3] if len(alt) >= 4 else ''
                    })

        response_payload = {
            'success': True,
            'question': question,
            'answer': answer,
            'matched_question': result.get('question', '') if result else '',
            'source': source,
            'preset': deepseek_meta.get('preset', preset),
            'preset_label': deepseek_meta.get('preset_label', '本地知识库' if use_local_only else ''),
            'model': deepseek_meta.get('model', ''),
            'alternatives': alternatives
        }
        if similarity is not None:
            response_payload['similarity'] = round(similarity, 3)
        if notice:
            response_payload['notice'] = notice

        return jsonify(response_payload)
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_robot/search', methods=['POST'])
@login_required
def api_qa_robot_search():
    """问答机器人搜索API - 返回多个相关结果"""
    if not QA_ROBOT_AVAILABLE:
        return jsonify({'success': False, 'error': '问答机器人模块不可用'}), 503
    
    try:
        data = request.get_json()
        query = data.get('query', '').strip()
        top_k = data.get('top_k', 5)
        threshold = data.get('threshold', 0.3)
        
        if not query:
            return jsonify({'success': False, 'error': '查询不能为空'}), 400
        
        # 搜索
        results = qa_robot_service.search(query, top_k=top_k, threshold=threshold)
        
        return jsonify({
            'success': True,
            'results': [
                {
                    'question': q,
                    'answer': a,
                    'similarity': round(s, 3)
                }
                for q, a, s in results
            ],
            'count': len(results)
        })
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_robot/history', methods=['GET'])
@login_required
def api_qa_robot_history():
    """获取问答机器人历史记录"""
    try:
        if 'user_id' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401
        
        limit = request.args.get('limit', 100, type=int)
        offset = request.args.get('offset', 0, type=int)
        
        history = get_qa_robot_history(session['user_id'], limit=limit, offset=offset)
        count = get_qa_robot_history_count(session['user_id'])
        
        return jsonify({
            'success': True,
            'history': history,
            'count': count
        })
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_robot/history/delete', methods=['POST'])
@login_required
def api_qa_robot_history_delete():
    """删除问答机器人历史记录"""
    try:
        if 'user_id' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401
        
        data = request.get_json()
        history_id = data.get('history_id', None)  # 如果为None，删除所有记录
        
        success, message = delete_qa_robot_history(session['user_id'], history_id)
        
        if success:
            return jsonify({'success': True, 'message': message})
        else:
            return jsonify({'success': False, 'error': message}), 400
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/translate', methods=['POST'])
@login_required
def api_translate():
    """翻译API"""
    if not TRANSLATOR_AVAILABLE:
        return jsonify({'success': False, 'error': '翻译模块不可用'}), 503
    
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({'success': False, 'error': '未提供文本内容'}), 400
        
        text = data['text'].strip()
        if not text:
            return jsonify({'success': False, 'error': '文本内容不能为空'}), 400
        
        # 获取目标语言，默认为中文
        target_lang = data.get('target_lang', 'zh')
        if target_lang not in ['zh', 'en', 'jp', 'ja']:
            target_lang = 'zh'
        
        # 执行翻译
        if target_lang == 'zh':
            translated_text = translate_to_chinese(text)
        elif target_lang in ['jp', 'ja']:
            translated_text = translate_to_japanese(text)
        else:
            translated_text = translate_to_english(text)
        
        return jsonify({
            'success': True,
            'original_text': text,
            'translated_text': translated_text,
            'target_lang': target_lang
        })
    
    except TranslationError as e:
        return jsonify({'success': False, 'error': str(e)}), 500
    except Exception as e:
        return jsonify({'success': False, 'error': f'服务器错误: {str(e)}'}), 500



# 语音助手API路由
@app.route('/api/voice/asr', methods=['POST'])
@login_required
def api_voice_asr():
    """语音识别API（支持自动识别普通话和英文）"""
    if not VOICE_ASSISTANT_AVAILABLE:
        return jsonify({'success': False, 'error': '语音助手模块不可用'}), 503
    
    try:
        if 'audio' not in request.files:
            return jsonify({'success': False, 'error': '未上传音频文件'}), 400
        
        file = request.files['audio']
        if file.filename == '':
            return jsonify({'success': False, 'error': '未选择文件'}), 400
        
        if not allowed_file(file.filename):
            return jsonify({'success': False, 'error': '不支持的文件格式，请上传 pcm, wav, amr 或 mp3 文件'}), 400
        
        # 获取是否自动检测语言的参数（默认为True）
        auto_detect = request.form.get('auto_detect', 'true').lower() == 'true'
        
        # 保存上传的文件
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        # 进行语音识别（自动识别语言）
        result = asr(filepath, auto_detect_language=auto_detect)
        
        # 删除临时文件
        try:
            os.remove(filepath)
        except:
            pass
        
        return jsonify({
            'success': True,
            'result': result
        })
    
    except ASRError as e:
        return jsonify({'success': False, 'error': f'识别错误: {str(e)}'}), 500
    except Exception as e:
        return jsonify({'success': False, 'error': f'服务器错误: {str(e)}'}), 500


@app.route('/api/voice/tts', methods=['POST'])
@login_required
def api_voice_tts():
    """语音生成API（支持选择中文或英文）"""
    if not VOICE_ASSISTANT_AVAILABLE:
        return jsonify({'success': False, 'error': '语音助手模块不可用'}), 503
    
    try:
        data = request.get_json()
        if not data or 'text' not in data:
            return jsonify({'success': False, 'error': '未提供文本内容'}), 400
        
        text = data['text'].strip()
        if not text:
            return jsonify({'success': False, 'error': '文本内容不能为空'}), 400
        
        # 获取语言参数，默认为中文
        language = data.get('language', 'zh')
        # 验证语言参数，确保是有效的值
        if language not in ['zh', 'en']:
            language = 'zh'  # 默认使用中文
        
        # 获取发音人参数，默认为0（度小美）
        speaker = data.get('speaker', 0)
        # 验证发音人参数，只允许0, 1, 3, 4
        if speaker not in [0, 1, 3, 4]:
            speaker = 0  # 默认使用度小美
        
        # 调试信息（可选，生产环境可以注释掉）
        print(f"[TTS API] 接收到的参数 - text长度: {len(text)}, language: {language}, speaker: {speaker}")
        
        # 先检查缓存
        cache_success, cached_file = get_tts_audio_cache(text, language)
        if cache_success and cached_file and os.path.exists(cached_file):
            # 使用缓存文件
            return send_file(
                cached_file,
                mimetype='audio/mpeg' if cached_file.endswith('.mp3') else 'audio/wav',
                as_attachment=True,
                download_name=os.path.basename(cached_file)
            )
        
        # 缓存不存在，生成新的音频文件
        # 使用文本和语言的哈希值生成唯一文件名
        import hashlib
        text_hash = hashlib.md5(text.encode('utf-8')).hexdigest()
        unique_filename = f"tts_{text_hash}_{language}.mp3"
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], unique_filename)
        
        # 切换到 voice_assistant 目录执行 tts（因为 tts.py 会在当前目录保存文件）
        original_dir = os.getcwd()
        try:
            os.chdir(voice_assistant_path)
            saved_file = tts(text, language=language, speaker=speaker)
        finally:
            os.chdir(original_dir)
        
        # 检查是否生成成功
        if saved_file == "error.txt":
            return jsonify({'success': False, 'error': '语音生成失败'}), 500
        
        # 处理文件路径
        if not os.path.isabs(saved_file):
            # 相对路径，检查 voice_assistant 目录
            file_in_va = os.path.join(voice_assistant_path, saved_file)
            if os.path.exists(file_in_va):
                # 移动到输出目录，使用唯一文件名
                try:
                    # 如果目标文件已存在，先删除
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    os.rename(file_in_va, output_path)
                    saved_file = output_path
                except OSError as e:
                    # 如果移动失败，尝试复制
                    import shutil
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    shutil.copy2(file_in_va, output_path)
                    saved_file = output_path
                    # 删除原文件
                    try:
                        os.remove(file_in_va)
                    except:
                        pass
            else:
                # 可能在输出目录中
                if os.path.exists(output_path):
                    saved_file = output_path
                else:
                    return jsonify({'success': False, 'error': '生成的音频文件未找到'}), 500
        
        # 保存到缓存
        save_tts_audio_cache(text, language, saved_file)
        
        # 返回音频文件
        return send_file(
            saved_file,
            mimetype='audio/mpeg' if saved_file.endswith('.mp3') else 'audio/wav',
            as_attachment=True,
            download_name=os.path.basename(saved_file)
        )
    
    except Exception as e:
        return jsonify({'success': False, 'error': f'服务器错误: {str(e)}'}), 500


@app.route('/admin')
@login_required
@admin_required
def admin_panel():
    """管理员后台页面"""
    return render_template('admin.html', username=session.get('username'))


@app.route('/admin/login', methods=['GET', 'POST'])
def admin_login():
    """管理员登录页面"""
    if request.method == 'POST':
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        
        if not username or not password:
            return jsonify({'success': False, 'error': '用户名和密码不能为空'}), 400
        
        # 验证用户
        success, message, user_data = verify_user(username, password)
        
        if not success:
            return jsonify({'success': False, 'error': message}), 401
        
        # 检查是否为管理员
        if user_data.get('role') != 'admin':
            return jsonify({'success': False, 'error': '该账号不是管理员'}), 403
        
        # 登录成功，设置session
        session['username'] = username
        session['user_id'] = user_data['id']
        session['role'] = 'admin'
        
        return jsonify({
            'success': True, 
            'message': '管理员登录成功', 
            'user': user_data
        })
    
    # GET请求，返回管理员登录页面
    if 'username' in session and session.get('role') == 'admin':
        return redirect(url_for('admin_panel'))
    return render_template('admin_login.html')


@app.route('/api/user_info')
@login_required
def user_info():
    """获取当前用户信息"""
    username = session.get('username')
    user_data = get_user_info(username)
    
    if user_data:
        return jsonify({'success': True, 'user': user_data})
    return jsonify({'success': False, 'error': '用户不存在'}), 404


@app.route('/api/admin/users', methods=['GET', 'POST'])
@login_required
@admin_required
def admin_users():
    """获取所有用户列表或创建新用户（管理员）"""
    if request.method == 'GET':
        users = get_all_users()
        return jsonify({'success': True, 'users': users})
    
    # POST: 创建新用户
    try:
        data = request.get_json()
        username = data.get('username', '').strip()
        password = data.get('password', '').strip()
        email = data.get('email', '').strip() or None
        role = data.get('role', 'user').strip()
        
        # 验证输入
        if not username or not password:
            return jsonify({'success': False, 'error': '用户名和密码不能为空'}), 400
        
        if len(username) < 3:
            return jsonify({'success': False, 'error': '用户名至少需要3个字符'}), 400
        
        if len(password) < 6:
            return jsonify({'success': False, 'error': '密码至少需要6个字符'}), 400
        
        if role not in ['user', 'admin']:
            return jsonify({'success': False, 'error': '无效的角色'}), 400
        
        # 创建用户
        success, message = register_user(username, password, email, role=role)
        
        if success:
            return jsonify({'success': True, 'message': message})
        else:
            status_code = 409 if '已存在' in message else 400
            return jsonify({'success': False, 'error': message}), status_code
    
    except Exception as e:
        return jsonify({'success': False, 'error': f'服务器错误: {str(e)}'}), 500


@app.route('/api/admin/user/<int:user_id>/role', methods=['PUT'])
@login_required
@admin_required
def admin_update_role(user_id):
    """更新用户角色（管理员）"""
    data = request.get_json()
    new_role = data.get('role', '').strip()
    
    if new_role not in ['user', 'admin']:
        return jsonify({'success': False, 'error': '无效的角色'}), 400
    
    success, message = update_user_role(user_id, new_role)
    
    if success:
        return jsonify({'success': True, 'message': message})
    return jsonify({'success': False, 'error': message}), 400


@app.route('/api/admin/tts/cache/delete', methods=['POST'])
@login_required
@admin_required
def admin_delete_tts_cache():
    """删除所有TTS音频缓存文件（管理员）"""
    try:
        data = request.get_json() or {}
        delete_db_records = data.get('delete_db_records', True)
        
        success, message, deleted_count = delete_all_tts_audio_files(delete_db_records=delete_db_records)
        
        if success:
            return jsonify({
                'success': True,
                'message': message,
                'deleted_count': deleted_count
            })
        else:
            return jsonify({'success': False, 'error': message}), 500
    
    except Exception as e:
        return jsonify({'success': False, 'error': f'服务器错误: {str(e)}'}), 500


@app.route('/api/admin/user/<int:user_id>', methods=['DELETE'])
@login_required
@admin_required
def admin_delete_user(user_id):
    """删除用户（管理员）"""
    current_user_id = session.get('user_id')
    success, message = delete_user(user_id, current_user_id)
    
    if success:
        return jsonify({'success': True, 'message': message})
    return jsonify({'success': False, 'error': message}), 400


# QA讨论区API路由
@app.route('/api/qa_discussion/submit', methods=['POST'])
@login_required
def api_qa_discussion_submit():
    """提交讨论内容"""
    try:
        if 'user_id' not in session or 'username' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401
        
        data = request.get_json()
        content = data.get('content', '').strip()
        
        if not content:
            return jsonify({'success': False, 'error': '内容不能为空'}), 400
        
        if len(content) > 5000:
            return jsonify({'success': False, 'error': '内容过长，最多5000字'}), 400
        
        success, message, discussion_id = save_qa_discussion(
            session['user_id'],
            session['username'],
            content
        )
        
        if success:
            return jsonify({
                'success': True,
                'message': message,
                'discussion_id': discussion_id
            })
        else:
            return jsonify({'success': False, 'error': message}), 400
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/list', methods=['GET'])
@login_required
def api_qa_discussion_list():
    """获取讨论列表"""
    try:
        limit = request.args.get('limit', 50, type=int)
        offset = request.args.get('offset', 0, type=int)
        
        discussions = get_qa_discussions(limit=limit, offset=offset)
        total = get_qa_discussion_count()
        
        # 为每个讨论获取评论数量
        for disc in discussions:
            comments = get_qa_comments(disc['id'])
            disc['comment_count'] = len(comments)
        
        return jsonify({
            'success': True,
            'discussions': discussions,
            'total': total
        })
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>', methods=['DELETE'])
@login_required
def api_qa_discussion_delete(discussion_id):
    """删除讨论"""
    try:
        if 'user_id' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401
        
        is_admin = session.get('role') == 'admin'
        user_id = session.get('user_id')
        
        success, message = delete_qa_discussion(discussion_id, user_id, is_admin)
        
        if success:
            return jsonify({'success': True, 'message': message})
        else:
            return jsonify({'success': False, 'error': message}), 400
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/comment', methods=['POST'])
@login_required
def api_qa_discussion_comment(discussion_id):
    """提交评论"""
    try:
        if 'user_id' not in session or 'username' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401
        
        data = request.get_json()
        content = data.get('content', '').strip()
        
        if not content:
            return jsonify({'success': False, 'error': '评论内容不能为空'}), 400
        
        if len(content) > 1000:
            return jsonify({'success': False, 'error': '评论过长，最多1000字'}), 400
        
        success, message, comment_id = save_qa_comment(
            discussion_id,
            session['user_id'],
            session['username'],
            content
        )
        
        if success:
            return jsonify({
                'success': True,
                'message': message,
                'comment_id': comment_id
            })
        else:
            return jsonify({'success': False, 'error': message}), 400
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/comments', methods=['GET'])
@login_required
def api_qa_discussion_comments(discussion_id):
    """获取讨论的评论列表"""
    try:
        comments = get_qa_comments(discussion_id)
        return jsonify({
            'success': True,
            'comments': comments
        })
    
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/comment/<int:comment_id>', methods=['DELETE'])
@login_required
def api_qa_discussion_comment_delete(comment_id):
    """删除评论"""
    try:
        if 'user_id' not in session:
            return jsonify({'success': False, 'error': '未登录'}), 401

        is_admin = session.get('role') == 'admin'
        user_id = session.get('user_id')

        success, message = delete_qa_comment(comment_id, user_id, is_admin)

        if success:
            return jsonify({'success': True, 'message': message})
        else:
            return jsonify({'success': False, 'error': message}), 400

    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


# ============================================================
# 讨论区增强 API：列表(支持排序/搜索)、详情、点赞/点踩、收藏、个人主页
# ============================================================

@app.route('/api/qa_discussion/feed', methods=['GET'])
@login_required
def api_qa_discussion_feed():
    """讨论区 Feed：支持最新/最热/搜索"""
    try:
        sort = request.args.get('sort', 'latest')  # 'latest' | 'hot'
        limit = min(request.args.get('limit', 20, type=int), 100)
        offset = max(request.args.get('offset', 0, type=int), 0)
        keyword = request.args.get('q', '').strip() or None

        user_id = session.get('user_id')
        discussions = get_qa_discussions_extended(sort=sort, limit=limit, offset=offset, keyword=keyword)
        total = get_qa_discussion_count_extended(keyword=keyword)

        # 批量补充：当前用户对这些讨论的投票/收藏状态
        ids = [d['id'] for d in discussions]
        vote_map = get_user_votes(user_id, 'discussion', ids) if ids else {}
        fav_map = get_user_favorited_ids(user_id, ids) if ids else {}

        for d in discussions:
            d['user_vote'] = vote_map.get(d['id'])
            d['is_favorited'] = fav_map.get(d['id'], False)

        return jsonify({
            'success': True,
            'discussions': discussions,
            'total': total,
            'has_more': offset + len(discussions) < total,
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/detail', methods=['GET'])
@login_required
def api_qa_discussion_detail(discussion_id):
    """获取讨论详情（增加浏览量）"""
    try:
        user_id = session.get('user_id')
        disc = get_qa_discussion_detail(discussion_id, increment_view=True)
        if not disc:
            return jsonify({'success': False, 'error': '讨论不存在'}), 404
        # 当前用户的投票/收藏状态
        votes = get_user_votes(user_id, 'discussion', [discussion_id])
        disc['user_vote'] = votes.get(discussion_id)
        disc['is_favorited'] = is_favorited(user_id, discussion_id)
        return jsonify({'success': True, 'discussion': disc})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/vote', methods=['POST'])
@login_required
def api_qa_discussion_vote(discussion_id):
    """对讨论点赞/点踩/取消"""
    try:
        user_id = session.get('user_id')
        data = request.get_json() or {}
        vote_type = data.get('vote_type', '')  # 'like' | 'dislike' | 'cancel'
        if vote_type not in ('like', 'dislike', 'cancel'):
            return jsonify({'success': False, 'error': 'vote_type 必须为 like/dislike/cancel'}), 400
        ok, msg, lc, dc, user_vote = vote_qa_target(user_id, 'discussion', discussion_id, vote_type)
        return jsonify({
            'success': ok,
            'message': msg,
            'like_count': lc,
            'dislike_count': dc,
            'user_vote': user_vote,
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/comment/<int:comment_id>/vote', methods=['POST'])
@login_required
def api_qa_comment_vote(comment_id):
    """对评论点赞/点踩/取消"""
    try:
        user_id = session.get('user_id')
        data = request.get_json() or {}
        vote_type = data.get('vote_type', '')
        if vote_type not in ('like', 'dislike', 'cancel'):
            return jsonify({'success': False, 'error': 'vote_type 必须为 like/dislike/cancel'}), 400
        ok, msg, lc, dc, user_vote = vote_qa_target(user_id, 'comment', comment_id, vote_type)
        return jsonify({
            'success': ok,
            'message': msg,
            'like_count': lc,
            'dislike_count': dc,
            'user_vote': user_vote,
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/favorite', methods=['POST'])
@login_required
def api_qa_discussion_favorite(discussion_id):
    """收藏/取消收藏讨论"""
    try:
        user_id = session.get('user_id')
        ok, msg, is_fav = toggle_favorite(user_id, discussion_id)
        return jsonify({'success': ok, 'message': msg, 'is_favorited': is_fav})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/qa_discussion/<int:discussion_id>/comments_extended', methods=['GET'])
@login_required
def api_qa_discussion_comments_extended(discussion_id):
    """获取评论列表（带点赞、@提及）"""
    try:
        user_id = session.get('user_id')
        comments = get_qa_comments_extended(discussion_id, user_id=user_id)
        return jsonify({'success': True, 'comments': comments})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


# ---------- 个人主页 ----------

@app.route('/user/<username>')
@login_required
def user_profile_page(username):
    """个人主页"""
    return render_template(
        'voice_assistant.html',
        username=session.get('username'),
        user_id=session.get('user_id'),
        current_view='user_profile',
        profile_username=username,
    )


@app.route('/api/user/<username>/profile', methods=['GET'])
@login_required
def api_user_profile(username):
    """获取个人主页概要"""
    profile = get_user_profile(username)
    if not profile:
        return jsonify({'success': False, 'error': '用户不存在'}), 404
    return jsonify({'success': True, 'profile': profile})


@app.route('/api/user/send_email_code', methods=['POST'])
@login_required
def api_send_email_code():
    """发送邮箱验证码 (模拟)"""
    import random
    try:
        data = request.get_json() or {}
        email = data.get('email', '').strip()
        if not email:
            return jsonify({'success': False, 'error': '邮箱不能为空'}), 400
        
        # 简单校验邮箱格式
        if '@' not in email or '.' not in email:
            return jsonify({'success': False, 'error': '邮箱格式不正确'}), 400

        code = str(random.randint(100000, 999999))
        session['email_code'] = code
        session['email_code_target'] = email
        print(f"[TEST EMAIL CODE] Sent verification code {code} to {email}")
        
        return jsonify({
            'success': True,
            'code': code,
            'message': f'验证码已发送至 {email}（测试模式，您的验证码为 {code}）'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/user/send_phone_code', methods=['POST'])
@login_required
def api_send_phone_code():
    """发送手机验证码 (模拟)"""
    import random
    try:
        data = request.get_json() or {}
        phone = data.get('phone', '').strip()
        if not phone:
            return jsonify({'success': False, 'error': '手机号不能为空'}), 400
        
        if len(phone) < 11:
            return jsonify({'success': False, 'error': '手机号格式不正确'}), 400

        code = str(random.randint(100000, 999999))
        session['phone_code'] = code
        session['phone_code_target'] = phone
        print(f"[TEST PHONE CODE] Sent verification code {code} to {phone}")
        
        return jsonify({
            'success': True,
            'code': code,
            'message': f'验证码已发送至 {phone}（测试模式，您的验证码为 {code}）'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/feature_experience/generate_outline', methods=['POST'])
@login_required
def api_generate_outline():
    """AI生成PPT大纲接口"""
    try:
        data = request.get_json() or {}
        topic = data.get('topic', '').strip()
        
        if not topic:
            return jsonify({'success': False, 'error': '主题不能为空'}), 400
            
        # 尝试从环境变量加载 Key
        from dotenv import load_dotenv
        load_dotenv(os.path.join(os.path.dirname(os.path.abspath(__file__)), '.env'))
        
        api_key = os.environ.get("MOARK_API_KEY") or os.environ.get("GITEE_AI_API_KEY") or os.environ.get("DEEPSEEK_API_KEY")
        if not api_key:
            return jsonify({'success': False, 'error': '未配置 API Key，无法使用 AI 生成功能'}), 500
            
        if os.environ.get("MOARK_API_KEY") or os.environ.get("GITEE_AI_API_KEY"):
            base_url = "https://api.moark.com/v1"
            model = "Qwen3-32B"
            extra_headers = {"X-Failover-Enabled": "true"}
        else:
            base_url = os.environ.get("DEEPSEEK_BASE_URL", "https://api.deepseek.com")
            model = os.environ.get("DEEPSEEK_MODEL", "deepseek-v4-flash")
            extra_headers = {}
            
        prompt = (
            f"请为PPT课件主题【{topic}】编写一份结构清晰、内容丰富的大纲。\n"
            f"大纲需要分章节，每个章节下列出每张幻灯片的主题，并说明其核心要点。\n"
            f"大纲格式应简洁，直接采用 Markdown 格式，不要有任何其他解释性废话，直接返回大纲内容。"
        )
        
        payload = {
            "model": model,
            "messages": [
                {"role": "system", "content": "You are a professional assistant that designs slide deck outlines. Output only the outline directly."},
                {"role": "user", "content": prompt}
            ],
            "temperature": 0.7,
            "max_tokens": 1500,
            "stream": False
        }
        
        if model == "Qwen3-32B":
            payload["top_p"] = 0.7
            payload["frequency_penalty"] = 1
            payload["extra_body"] = {"top_k": 50}
            
        res = deepseek_service._post_chat_completions(payload, api_key=api_key, base_url=base_url, extra_headers=extra_headers)
        choices = res.get("choices") or []
        if choices:
            content = choices[0].get("message", {}).get("content", "").strip()
            if not content and choices[0].get("message", {}).get("reasoning_content"):
                content = choices[0].get("message", {}).get("reasoning_content", "").strip()
            return jsonify({'success': True, 'outline': content})
        else:
            return jsonify({'success': False, 'error': 'API 未返回大纲内容'})
            
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500



@app.route('/api/user/update_profile', methods=['POST'])
@login_required
def api_update_profile():
    """更新用户资料接口"""
    try:
        data = request.get_json() or {}
        new_username = data.get('username', '').strip()
        new_email = data.get('email', '').strip() or None
        email_code = data.get('email_code', '').strip()
        new_phone = data.get('phone', '').strip() or None
        phone_code = data.get('phone_code', '').strip()
        new_birthday = data.get('birthday', '').strip() or None
        
        user_id = session.get('user_id')
        current_username = session.get('username')
        
        if not new_username:
            return jsonify({'success': False, 'error': '用户名不能为空'}), 400
        if len(new_username) < 3:
            return jsonify({'success': False, 'error': '用户名至少需要3个字符'}), 400

        # 获取当前的用户信息以进行修改前的对比
        current_user = get_user_info(current_username)
        if not current_user:
            return jsonify({'success': False, 'error': '未找到当前登录的用户信息'}), 404
        
        # 1. 邮箱验证码校验
        # 如果新邮箱不为空且与当前邮箱不同，则必须校验验证码
        current_email = current_user.get('email')
        if new_email and new_email != current_email:
            saved_code = session.get('email_code')
            saved_target = session.get('email_code_target')
            if not email_code:
                return jsonify({'success': False, 'error': '需要输入邮箱验证码'}), 400
            if saved_target != new_email or saved_code != email_code:
                return jsonify({'success': False, 'error': '邮箱验证码不正确或目标邮箱不匹配'}), 400
        
        # 2. 手机验证码校验
        # 如果新手机号不为空且与当前手机号不同，则必须校验验证码
        current_phone = current_user.get('phone')
        if new_phone and new_phone != current_phone:
            saved_code = session.get('phone_code')
            saved_target = session.get('phone_code_target')
            if not phone_code:
                return jsonify({'success': False, 'error': '需要输入手机验证码'}), 400
            if saved_target != new_phone or saved_code != phone_code:
                return jsonify({'success': False, 'error': '手机验证码不正确或目标手机号不匹配'}), 400

        # 3. 生日格式校验
        if new_birthday:
            import datetime
            try:
                datetime.datetime.strptime(new_birthday, '%Y-%m-%d')
            except ValueError:
                return jsonify({'success': False, 'error': '生日格式不正确，应为 YYYY-MM-DD'}), 400

        # 4. 执行更新
        success, message = update_user_profile(
            user_id=user_id,
            new_username=new_username,
            new_email=new_email,
            new_phone=new_phone,
            new_birthday=new_birthday
        )
        
        if not success:
            return jsonify({'success': False, 'error': message}), 400

        # 5. 更新成功后同步更新 session 里的用户名
        session['username'] = new_username
        
        # 清除已使用的验证码
        session.pop('email_code', None)
        session.pop('email_code_target', None)
        session.pop('phone_code', None)
        session.pop('phone_code_target', None)
        
        return jsonify({
            'success': True,
            'message': '个人资料更新成功！'
        })
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/user/<username>/discussions', methods=['GET'])
@login_required
def api_user_discussions(username):
    """个人主页：发布的讨论"""
    try:
        limit = min(request.args.get('limit', 20, type=int), 100)
        offset = max(request.args.get('offset', 0, type=int), 0)
        viewer_id = session.get('user_id')
        profile = get_user_profile(username)
        if not profile:
            return jsonify({'success': False, 'error': '用户不存在'}), 404
        items = get_user_discussions(profile['user']['id'], limit=limit, offset=offset)
        ids = [i['id'] for i in items]
        fav_map = get_user_favorited_ids(viewer_id, ids) if ids else {}
        for it in items:
            it['is_favorited'] = fav_map.get(it['id'], False)
        return jsonify({'success': True, 'discussions': items})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/user/<username>/comments', methods=['GET'])
@login_required
def api_user_comments(username):
    """个人主页：发表的评论"""
    try:
        limit = min(request.args.get('limit', 20, type=int), 100)
        offset = max(request.args.get('offset', 0, type=int), 0)
        profile = get_user_profile(username)
        if not profile:
            return jsonify({'success': False, 'error': '用户不存在'}), 404
        items = get_user_comments(profile['user']['id'], limit=limit, offset=offset)
        return jsonify({'success': True, 'comments': items})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500


@app.route('/api/user/<username>/favorites', methods=['GET'])
@login_required
def api_user_favorites(username):
    """个人主页：收藏的讨论（仅本人可看）"""
    try:
        viewer = session.get('username')
        if viewer != username and session.get('role') != 'admin':
            return jsonify({'success': False, 'error': '无权查看他人收藏'}), 403
        limit = min(request.args.get('limit', 20, type=int), 100)
        offset = max(request.args.get('offset', 0, type=int), 0)
        profile = get_user_profile(username)
        if not profile:
            return jsonify({'success': False, 'error': '用户不存在'}), 404
        items = get_user_favorites(profile['user']['id'], limit=limit, offset=offset)
        for it in items:
            it['is_favorited'] = True
        return jsonify({'success': True, 'favorites': items})
    except Exception as e:
        return jsonify({'success': False, 'error': str(e)}), 500

# ==================== 功能体验 API 路由 ====================

@app.route('/api/feature_experience/generate_ppt', methods=['POST'])
@login_required
def api_generate_ppt():
    """AI生成PPT接口"""
    try:
        data = request.get_json() or {}
        topic = data.get('topic', '').strip()
        style = data.get('style', 'tech').strip()
        page_count = data.get('page_count', 5)
        
        if not topic:
            return jsonify({'success': False, 'error': '主题不能为空'}), 400
            
        # 尝试调用 DeepSeek
        if DEEPSEEK_AVAILABLE and deepseek_service and deepseek_service.is_configured():
            prompt = (
                f"你是一个专业的PPT大纲及讲解内容生成助手。请根据用户的主题：“{topic}”，幻灯片风格风格为：“{style}”（古典雅致/现代科技/简约学术/活泼卡通），"
                f"生成共{page_count}页的PPT幻灯片内容。确保每页的讲解配音自然口语化，适合AI数字人播报。\n"
                f"你必须返回一个符合以下JSON数组格式的文本，不要包含任何Markdown标记（如```json），直接以 [ 开头，以 ] 结尾，以便于程序直接解析。\n"
                f"格式要求：\n"
                f"[\n"
                f"  {{\n"
                f"    \"title\": \"第1页的标题\",\n"
                f"    \"points\": [\"核心要点1\", \"核心要点2\", \"核心要点3\"],\n"
                f"    \"narration\": \"第一页幻灯片的讲解内容，数字人授课配音文本，要求通俗易懂，口语化，长度在60-120字左右。\"\n"
                f"  }},\n"
                f"  ...\n"
                f"]"
            )
            payload = {
                "model": deepseek_service.model,
                "messages": [
                    {"role": "system", "content": "You are a professional assistant that generates PPT slides in JSON format. Return only valid JSON array. Do not include markdown tags."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.5,
                "max_tokens": 3000,
                "stream": False
            }
            try:
                res = deepseek_service._post_chat_completions(payload)
                choices = res.get("choices") or []
                if choices:
                    content = choices[0].get("message", {}).get("content", "").strip()
                    if content.startswith("```"):
                        lines = content.split("\n")
                        if lines[0].startswith("```"):
                            lines = lines[1:]
                        if lines[-1].startswith("```"):
                            lines = lines[:-1]
                        content = "\n".join(lines).strip()
                    slides = json.loads(content)
                    return jsonify({'success': True, 'slides': slides, 'source': 'AI (DeepSeek)'})
            except Exception as e:
                print(f"[API_GENERATE_PPT] DeepSeek API 调用失败: {e}，将回退到本地模板引擎生成。")

        # 本地生成逻辑（Fallback）
        slides = []
        topic_lower = topic.lower()
        
        # 1. 匹配古诗词/中国文学/历史/传统文化
        if any(kw in topic_lower for kw in ['古诗', '诗词', '文学', '历史', '传统文化', '中国', '语文', '古代']):
            slides = [
                {
                    "title": f"《{topic}》的背景与文化渊源",
                    "points": ["历史朝代的兴衰与文学流派演变", "作者个人生平际遇与情感底色", "特定时代背景下的社会风貌"],
                    "narration": f"同学们好！今天我们来学习关于《{topic}》的精彩内容。首先，我们需要了解它的背景与文化渊源。每个优秀的文化作品，都深深地烙印着它所处的时代印记、社会风貌以及作者波澜壮阔的人生旅程。让我们穿越时光，回到那段历史之中。"
                },
                {
                    "title": "作品的核心意象与美学价值",
                    "points": ["经典自然意象的寄托与象征意蕴", "情景交融、虚实结合的艺术手法", "含蓄隽永、意境深远的独特美学风格"],
                    "narration": f"接下来，我们深入探讨《{topic}》中的美学价值。文学作品往往通过具体的意象来表达情感，比如月亮、春风、流水。作者善于把主观情感融入到客观景物中，达到情景交融的最高境界。我们品读时，不仅要赏析文字，更要感悟那份悠远的意境。"
                },
                {
                    "title": "名篇名句赏析与语言艺术",
                    "points": ["关键字词的精妙锤炼与多重释义", "句式结构对情感跌宕起伏的推波助澜", "声律和谐、琅琅上口的音韵之美"],
                    "narration": "学习古典文学，最美妙的莫过于赏析名篇名句。这些字句经过了千锤百炼，极其精妙。让我们一起来看这几个关键句，它们字句凝练，音韵和谐，读起来琅琅上口，在情感的高潮处推波助澜，表现力极强。"
                },
                {
                    "title": "对后世的深远影响与现代启示",
                    "points": ["在后世文学创作中的传承与发扬", "跨越时空阻隔的人性共鸣与情感投射", "现代语境下的创造性转化与创新发展"],
                    "narration": f"最后，我们来思考《{topic}》对后世的深远影响。优秀的文学和历史遗产是跨越时空的，它不仅深深影响了后世的创作，在今天也依然能引起我们的强烈共鸣。我们在现代语境下，应该如何去继承和发扬这些珍贵的文化财富呢？这值得我们每个人深思。"
                }
            ]
        # 2. 匹配人工智能/科技/教育信息化
        elif any(kw in topic_lower for kw in ['人工智能', '科技', 'ai', '教育', '智慧', '学校', '计算机', '技术']):
            slides = [
                {
                    "title": f"{topic} 的概念界定与兴起背景",
                    "points": ["前沿技术的迭代发展与政策环境红利", "传统教学模式面临的关键痛点与变革诉求", "人机协同教育新范式的概念演进"],
                    "narration": f"大家好，今天我们的课程主题是《{topic}》。首先，我们来看一下它的概念界定与兴起背景。近年来，随着大模型、算力等底层技术的爆发，以及国家教育数字化战略的推动，传统教育正经历深刻重塑。如何用前沿科技解决传统教学的痛点，是我们探讨的核心。"
                },
                {
                    "title": "核心应用场景与关键技术架构",
                    "points": ["多模态互动教学与智能学习资源生成", "个性化自适应学习路径精准规划", "全过程学习行为多维分析与智能评估"],
                    "narration": "那么，它有哪些核心的应用场景呢？首先是教学资源的生成，比如我们现在正在体验的数字人授课和PPT自动生成。其次是自适应学习，通过分析学生的行为，规划最适合的专属路线。最后是智能的学情评估，实现真正的因材施教。"
                },
                {
                    "title": "面临的现实挑战与伦理考量",
                    "points": ["数字鸿沟引发教育公平性的全新审视", "教学数据隐私保护与信息系统安全底线", "数字技术对教师传统角色定位 of 重塑"],
                    "narration": "虽然科技带来了极大的便利，但我们也必须保持清醒，看到面临的挑战。比如，落后地区是否能公平享受到这些技术？学生的隐私数据该如何严格保密？更重要的是，当AI能够做很多事情时，我们教师传统的工作角色该如何调整和定位？"
                },
                {
                    "title": "未来发展趋势与融合共生图景",
                    "points": ["大模型技术在细分教育垂直领域的深度融合", "人机双能驱动教学体系的全面建立", "回归教育本质——关注学生核心素养培养"],
                    "narration": "最后，展望未来，技术必将与教育各环节进行更深度的融合，构建一个教师与AI双能驱动、融合共生的智能教育新图景。无论技术怎么变，教育的本质——点燃智慧、润泽心灵是永远不会改变的。让我们拥抱科技，共同见证教育的美好未来！"
                }
            ]
        # 3. 其他常规主题
        else:
            slides = [
                {
                    "title": f"关于《{topic}》的引入与概览",
                    "points": ["探索该领域的研究初衷与核心价值", "行业发展的重要阶段与历史沿革", "课程内容的核心架构与知识版图"],
                    "narration": f"各位学员好！今天我将带大家一起探索关于《{topic}》的主题。本课程旨在从多维视角分析该领域的核心框架、演进历史以及未来的关键方向。让我们首先从它的起源和研究初衷开始，逐步揭开它的神秘面纱。"
                },
                {
                    "title": "核心内容解析与要点细化",
                    "points": ["基本理论架构与底层逻辑支撑", "不同流派或观点的碰撞与技术交融", "实现该领域突破的核心要素分析"],
                    "narration": f"接下来，我们切入《{topic}》的核心内容进行解析。在这部分，我们将重点讨论其底层的理论逻辑与核心支撑要素。通过对比不同的观点和应用流派，我们将理清该领域在实践中取得重大突破的深层原因。"
                },
                {
                    "title": "应用实例与典型案例分析",
                    "points": ["国内外典型成功案例深度剖析", "实施过程中的关键路径与避坑指南", "从案例中提炼的普适性规律与实操经验"],
                    "narration": "理论的生命力在于实践。让我们来看几个具有代表性的成功案例。通过深度剖析这几个案例在不同发展阶段所做出的关键决策，我们可以提炼出许多可以直接借鉴的实操经验，避免在实际推进中少走弯路。"
                },
                {
                    "title": "总结回顾与前瞻性展望",
                    "points": ["本堂课核心概念的串联与知识收敛", "该领域下一步发展的核心驱动力", "给每位学习者的个人成长与实践建议"],
                    "narration": f"最后，我们对今天所学关于《{topic}》的内容做个总结。我们梳理了它的理论框架与案例，展望了它接下来的爆发点。希望大家在课后能够结合自己的工作和学习进行实践。谢谢大家，我们下期课程再见！"
                }
            ]
            
        # 如果请求的页数较少，进行裁剪；如果较多，我们复制填充或者按模板补充
        if len(slides) > page_count:
            slides = slides[:page_count]
        elif len(slides) < page_count:
            # 简单补充
            for i in range(len(slides), page_count):
                slides.append({
                    "title": f"关于《{topic}》的深入探究 (页 {i+1})",
                    "points": ["探索未知维度的补充概念与支撑材料", "关键论点的延伸讨论与论据丰富", "实践层面的更多操作规范与应用指导"],
                    "narration": f"在这里，我们对《{topic}》的相关内容进行进一步的延伸探讨。在本页，我们将细化补充论点，确保我们的知识体系更加健全，帮助大家在实际操作中具备更全面的视野。"
                })
                
        return jsonify({'success': True, 'slides': slides, 'source': 'Local Template'})
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'PPT生成失败: {str(e)}'}), 500


@app.route('/api/feature_experience/correct_homework', methods=['POST'])
@login_required
def api_correct_homework():
    """AI批改作业接口"""
    try:
        if request.is_json:
            data = request.get_json() or {}
            homework_text = data.get('homework_text', '').strip()
            subject = data.get('subject', 'chinese').strip()
        else:
            homework_text = request.form.get('homework_text', '').strip()
            subject = request.form.get('subject', 'chinese').strip()
            
        if not homework_text:
            return jsonify({'success': False, 'error': '作业内容不能为空'}), 400
            
        # 尝试调用 DeepSeek
        if DEEPSEEK_AVAILABLE and deepseek_service and deepseek_service.is_configured():
            prompt = (
                f"你是一个资深的教育导师。请对以下学生的“{subject}”科目作业进行精细化的批改。反馈需要温和鼓励但指出关键错误。\n"
                f"学生作业内容：\n"
                f"\"\"\"\n{homework_text}\n\"\"\"\n\n"
                f"请分析作业中的语法错误、逻辑错误、计算错误或笔误，并给出改进意见。\n"
                f"你必须返回一个符合以下JSON格式的纯文本，不要包含任何Markdown标记（如```json），以便于程序解析：\n"
                f"{{\n"
                f"  \"score\": 评分(1-100的整数),\n"
                f"  \"overall_comment\": \"总评，夸奖优点并委婉指出不足，富有启发性(100-150字左右)\",\n"
                f"  \"corrections\": [\n"
                f"    {{\n"
                f"      \"original\": \"原错词/句/算式\",\n"
                f"      \"corrected\": \"修改后的正确内容\",\n"
                f"      \"reason\": \"详细的纠错原因与知识点解析\"\n"
                f"    }},\n"
                f"    ...\n"
                f"  ],\n"
                f"  \"suggestions\": [\n"
                f"    \"后续具体的复习与练习建议1\",\n"
                f"    \"后续具体的复习与练习建议2\"\n"
                f"  ]\n"
                f"}}"
            )
            payload = {
                "model": deepseek_service.model,
                "messages": [
                    {"role": "system", "content": "You are a professional educational teacher grading homework. Return only a valid JSON object. Do not include markdown formatting."},
                    {"role": "user", "content": prompt}
                ],
                "temperature": 0.4,
                "max_tokens": 2000,
                "stream": False
            }
            try:
                res = deepseek_service._post_chat_completions(payload)
                choices = res.get("choices") or []
                if choices:
                    content = choices[0].get("message", {}).get("content", "").strip()
                    if content.startswith("```"):
                        lines = content.split("\n")
                        if lines[0].startswith("```"):
                            lines = lines[1:]
                        if lines[-1].startswith("```"):
                            lines = lines[:-1]
                        content = "\n".join(lines).strip()
                    result = json.loads(content)
                    return jsonify({'success': True, 'report': result, 'source': 'AI (DeepSeek)'})
            except Exception as e:
                print(f"[API_CORRECT_HOMEWORK] DeepSeek API 调用失败: {e}，将回退到本地启发式批改。")

        # 本地批改逻辑 (Fallback)
        report = {
            "score": 85,
            "overall_comment": "",
            "corrections": [],
            "suggestions": []
        }
        
        text_lower = homework_text.lower()
        if subject == 'math':
            report["overall_comment"] = "计算过程比较清晰，书写也规范，但是在部分基础算式和公式套用上出现了笔误。建议日常练习时多加复查，养成细心的好习惯。"
            if '1+1=3' in homework_text.replace(' ', ''):
                report["corrections"].append({
                    "original": "1 + 1 = 3",
                    "corrected": "1 + 1 = 2",
                    "reason": "最基础的加法计算失误，请集中注意力进行运算。"
                })
            if '3x=6' in homework_text.replace(' ', '') and 'x=3' in homework_text.replace(' ', ''):
                report["corrections"].append({
                    "original": "3x = 6 推出 x = 3",
                    "corrected": "x = 2",
                    "reason": "方程两边应同时除以系数3，得到 x = 6/3 = 2。你可能误用减法去移项了。"
                })
            
            if not report["corrections"]:
                report["corrections"] = [
                    {
                        "original": "计算中间步骤的单位漏写",
                        "corrected": "带上对应单位（如：cm² 或 元）",
                        "reason": "应用题的计算过程中，关键步骤建议保留单位，最终结果必须注明正确单位，以免扣分。"
                    },
                    {
                        "original": "除法运算中余数写错",
                        "corrected": "仔细验算：被除数 = 除数 × 商 + 余数",
                        "reason": "除法竖式计算中，减法借位出现了偏差，导致余数大于除数，请牢记余数必须小于除数的定理。"
                    }
                ]
            report["suggestions"] = [
                "准备一本错题本，把平时因为马虎写错的题目收集起来，考试前看一遍。",
                "做完题后，使用代入法或者逆运算法进行一轮快速的口算验算。",
                "每天坚持进行5道口算或心算基础练习，提升对数字的敏感度。"
            ]
            report["score"] = 82
            
        elif subject == 'english':
            report["overall_comment"] = "An interesting writing with rich ideas! Your vocabulary is good, but there are a few grammatical errors related to subject-verb agreement and preposition usage. Keep writing!"
            if 'i is' in text_lower or 'he are' in text_lower or 'she are' in text_lower:
                report["corrections"].append({
                    "original": "I is / He are / She are",
                    "corrected": "I am / He is / She is",
                    "reason": "Be动词需与主语人称保持一致。第一人称单数I用am，第三人称单数he/she/it用is。"
                })
            if 'good at english' not in text_lower and 'good in english' in text_lower:
                report["corrections"].append({
                    "original": "good in English",
                    "corrected": "good at English",
                    "reason": "固定搭配。表示擅长于某项学科或技能用 be good at，后接名词或动名词形式。"
                })
            if 'last year i go' in text_lower or 'yesterday i buy' in text_lower:
                report["corrections"].append({
                    "original": "last year I go / yesterday I buy",
                    "corrected": "last year I went / yesterday I bought",
                    "reason": "时态错误。时间状语表示过去时间（last year, yesterday等）时，谓语动词应使用过去时态。"
                })
                
            if not report["corrections"]:
                report["corrections"] = [
                    {
                        "original": "The student show a lot of interest...",
                        "corrected": "The student shows a lot of interest...",
                        "reason": "主谓一致。当主语是第三人称单数时，一般现在时的谓语动词要加上词尾 -s 或 -es。"
                    },
                    {
                        "original": "in sunday morning",
                        "corrected": "on Sunday morning",
                        "reason": "介词搭配。表示在具体某一天或某一天的上午、下午、晚上时，前面要用介词 on，首字母 Sunday 需大写。"
                    }
                ]
            report["suggestions"] = [
                "Pay attention to verb tenses when describing past events.",
                "Review the rules of third-person singular verbs in present simple tense.",
                "Read english articles aloud for 10 minutes daily to build natural grammar intuition."
            ]
            report["score"] = 88
            
        else: # Chinese
            report["overall_comment"] = "文章构思新颖，情感表达真挚，细节描写非常生动。但是，在个别段落中存在着错别字，以及词语搭配不够恰当的情况，精雕细琢后会是一篇佳作！"
            if '得时候' in homework_text:
                report["corrections"].append({
                    "original": "得时候",
                    "corrected": "的时候",
                    "reason": "结构助词使用错误。表示时间的“...的时候”应该使用“的”字，而不是表示程度或补语的“得”。"
                })
            if '克苦' in homework_text or '客苦' in homework_text:
                report["corrections"].append({
                    "original": "克苦 / 客苦",
                    "corrected": "刻苦",
                    "reason": "字形写错。“刻苦”的“刻”是雕刻的刻，意思是像雕刻一样深刻用力，形容非常吃苦、努力。"
                })
            
            if not report["corrections"]:
                report["corrections"] = [
                    {
                        "original": "他那红润的脸蛋上露出了欣慰的笑容和开心的笑声",
                        "corrected": "他那红润的脸蛋上露出了欣慰的笑容，传传来开心的笑声",
                        "reason": "动宾搭配不当。笑容可以“露出”，但笑声只能“传来”或“听到”，不能与“露出”搭配使用。"
                    },
                    {
                        "original": "通过这次活动，使我深刻懂得了团队合作的重要性",
                        "corrected": "这次活动使我深刻懂得了... / 通过这次活动，我深刻懂得了...",
                        "reason": "成分残缺（主语缺失）。“通过...”和“使...”连用，导致句子缺少主语。可去掉“通过”或去掉“使”。"
                    }
                ]
            report["suggestions"] = [
                "在写作完成后，养成默读一到两遍的习惯，自主找出语气不顺的病句和错别字。",
                "多积累好词好句，注意动词和名词之间的搭配合理性，做到准确表达。",
                "重点温习多音字和形近字，把容易写错的字单独抄写三遍加深记忆。"
            ]
            report["score"] = 86
            
        return jsonify({'success': True, 'report': report, 'source': 'Local Rules Engine'})
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'批改失败: {str(e)}'}), 500


@app.route('/api/feature_experience/upload_ppt', methods=['POST'])
@login_required
def api_upload_ppt():
    """上传并解析 PPTX 课件"""
    if 'file' not in request.files:
        return jsonify({'success': False, 'error': '未选择文件'}), 400
        
    file = request.files['file']
    if file.filename == '':
        return jsonify({'success': False, 'error': '未选择文件'}), 400
        
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({'success': False, 'error': '目前仅支持解析 .pptx 格式的幻灯片文件'}), 400
        
    try:
        # 保存临时文件
        filename = secure_filename(file.filename)
        filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(filepath)
        
        slides = []
        
        # 使用 python-pptx 解析
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            for i, slide in enumerate(prs.slides):
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                
                points = []
                for shape in slide.shapes:
                    if shape.has_text_frame and (not slide.shapes.title or shape != slide.shapes.title):
                        for paragraph in shape.text_frame.paragraphs:
                            txt = paragraph.text.strip()
                            if txt and txt not in points:
                                points.append(txt)
                
                # 如果没有标题，尝试使用第一个文本点作为标题
                if not title and points:
                    title = points.pop(0)
                
                if not title:
                    title = f"幻灯片 {i+1}"
                
                # 限制要点字数和个数，过滤无意义空白
                points = [p for p in points if len(p) > 1][:4]
                
                # 自动为数字人生成讲解词
                if points:
                    points_summary = "，".join([p[:20] for p in points[:3]])
                    narration = f"现在我们来看这一页，它的核心内容是：{title}。主要包括以下几个要点：{points_summary}。请大家仔细对照屏幕上的内容进行理解。"
                else:
                    narration = f"大家请看这一页幻灯片，它的主题是：{title}。这里展示了相关的重要概念，请大家进行阅读和思考。"
                    
                slides.append({
                    "title": title,
                    "points": points if points else ["图片、图表或多媒体展示页"],
                    "narration": narration
                })
        except ImportError:
            try:
                os.remove(filepath)
            except:
                pass
            return jsonify({'success': False, 'error': '服务器未配置 python-pptx 解析环境，请联系管理员运行 `pip install python-pptx`。'}), 500
            
        # 删除临时文件
        try:
            os.remove(filepath)
        except:
            pass
            
        if not slides:
            return jsonify({'success': False, 'error': '未能在该幻灯片文件中解析出任何有效的文字页面。'}), 400
            
        return jsonify({'success': True, 'slides': slides})
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        return jsonify({'success': False, 'error': f'幻灯片解析失败: {str(e)}'}), 500


if __name__ == '__main__':
    print("="*50)
    print("智能问答系统启动中...")
    print("="*50)
    
    # 初始化数据库
    try:
        print("正在初始化数据库...")
        init_database()
        print("数据库初始化成功")
    except Exception as e:
        print(f"数据库初始化失败: {e}")
        print("请检查MySQL服务是否运行，以及config.py中的配置是否正确")
        exit(1)
    
    print("="*60)
    print("🚀 服务器启动成功！")
    print("="*60)
    print("📌 访问地址:")
    print("   本地访问: http://127.0.0.1:5001")
    print("   局域网访问: http://0.0.0.0:5001")
    print("   端口: 5001")
    print("="*60)
    print("💡 提示: 按 Ctrl+C 停止服务器")
    print("="*60)
    print()
    app.run(debug=True, host='0.0.0.0', port=5001)
